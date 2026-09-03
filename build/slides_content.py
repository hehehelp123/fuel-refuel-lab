# -*- coding: utf-8 -*-
"""Содержимое слайдов. Подключается из build_slides_itmo.py."""
import os

import pandas as pd
from pptx.util import Inches, Pt

CSV = "C:/Users/79101/Downloads/Лаба_заправки/data/fuel_5min.csv"


def build(ctx):
    add = ctx["add"]
    add_two_pics = ctx["add_two_pics"]
    stat = ctx["stat"]
    new_slide = ctx["new_slide"]
    fill = ctx["fill"]
    stamp_page_no = ctx["stamp_page_no"]
    L = ctx["L_PURPLE"]
    PURPLE, INK, WHITE = ctx["PURPLE"], ctx["INK"], ctx["WHITE"]
    prs = ctx["prs"]

    # ------------------------------------------------------------ таблица
    def add_table(slide, rows, x, y, w, col_w=None, size=12, head_size=12):
        n_rows, n_cols = len(rows), len(rows[0])
        h = 0.34 * n_rows
        shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                       Inches(w), Inches(h))
        table = shape.table
        table.first_row = True
        table.horz_banding = False
        if col_w:
            for i, cw in enumerate(col_w):
                table.columns[i].width = Inches(cw)
        for r, row in enumerate(rows):
            table.rows[r].height = Inches(0.34)
            for c, text in enumerate(row):
                cell = table.cell(r, c)
                cell.text = ""
                cell.margin_left = Inches(0.08)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                para = cell.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = str(text)
                run.font.size = Pt(head_size if r == 0 else size)
                run.font.bold = (r == 0)
                if r == 0:
                    run.font.color.rgb = WHITE
        return shape

    def text_slide_with_table(title, rows, note=None, col_w=None, size=12):
        slide = new_slide(L, title)
        for ph in list(slide.placeholders):
            if ph.placeholder_format.idx != 0:
                ph._element.getparent().remove(ph._element)
        y = 1.55
        add_table(slide, rows, 0.85, y, 8.3, col_w=col_w, size=size)
        if note:
            box = slide.shapes.add_textbox(Inches(0.85), Inches(y + 0.34 * len(rows) + 0.2),
                                           Inches(8.3), Inches(0.6))
            box.text_frame.word_wrap = True
            fill(box.text_frame, [{"t": note, "size": 14, "color": INK}])
        stamp_page_no(slide)
        return slide

    # ================================================================= задача
    add(L, "Задача", [
        {"t": "Датчик уровня топлива шлёт показания раз в несколько секунд. "
              "Надо находить моменты заправок.", "size": 15},
    ] + stat("1249", "заправок размечено за полгода на 50 машинах"),
        body_box=(5.7, 1.7, 3.6, 2.6),
        picture="s_events_split.png", pic_box=(0.7, 2.0, 4.7))

    # ------------------------------------------------------- что приходит
    text_slide_with_table(
        "Что приходит с машины",
        [["поле", "что это"],
         ["уровень топлива", "в единицах датчика, не в литрах"],
         ["зажигание", "0 или 1; при нуле датчик молчит"],
         ["скорость", "отличает стоянку от движения"],
         ["время", "сообщения приходят неравномерно"]],
        note="Перевод в литры — по тарировочной таблице своей машины. "
             "Без неё одна заправка на разных машинах даёт разную амплитуду.",
        col_w=[2.5, 5.8])

    # ------------------------------------------------------------- порог
    add(L, "Порог на приросте не работает", [
        {"t": "Низкий порог ловит колебания уровня, высокий пропускает "
              "заправки. Середины нет.", "size": 15},
    ] + stat("0.37", "лучший событийный F1 порога"),
        body_box=(5.9, 1.8, 3.4, 2.4),
        picture="s_threshold.png", pic_box=(0.7, 1.9, 5.0))

    # --------------------------------------------------- две ситуации
    add_two_pics("Зажигание делит задачу надвое", L,
                 ("Двигатель выключен", "s_ign_off.png",
                  "Датчик молчит всю заправку: 8 л до остановки, 60 л после запуска. 711 событий"),
                 ("Двигатель включён", "s_event.png",
                  "Машина стоит, двигатель работает, датчик пишет рост. 538 событий"))

    # ----------------------------------------------------------- датасет
    df = pd.read_csv(CSV, parse_dates=["ts"])
    veh = df[df.label == 1].vehicle.value_counts().index[0]
    g = df[df.vehicle == veh].reset_index(drop=True)
    i = g.index[g.label == 1][len(g.index[g.label == 1]) // 2]
    sample = g.iloc[max(0, i - 2):i + 4]
    rows = [["vehicle", "ts", "fuel", "ign", "speed", "label"]]
    for _, r in sample.iterrows():
        rows.append([r.vehicle, f"{r.ts:%d.%m %H:%M}", f"{r.fuel:.1f}",
                     f"{r.ign:.1f}", f"{r.speed:.0f}", int(r.label)])
    text_slide_with_table(
        "Датасет: 85 228 строк",
        rows,
        note="Одна строка — 5 минут одной машины. Разбиение по машинам: "
             "36 обучение, 9 валидация, 5 тест. Классы: 0 — ничего, 1 — заправка "
             "с включённым двигателем, 2 — с выключенным.",
        col_w=[1.2, 1.8, 1.3, 1.1, 1.3, 1.6], size=11)

    # --------------------------------------------------------- дисбаланс
    add(L, "Целевой класс редкий", stat("1.8%", "интервалов — заправка с включённым двигателем") + [
        {"t": "Модель «всегда ноль» даёт 98% accuracy", "size": 14, "space": 12},
        {"t": "Спасают веса классов в функции потерь", "size": 14},
    ], body_box=(6.0, 1.7, 3.3, 2.6),
        picture="s_balance.png", pic_box=(0.65, 1.9, 5.1))

    # ------------------------------------------------------ признаки и окна
    add(L, "Признаки и окна", [
        {"t": "Признаки интервала: уровень и четыре разности — на 1, 2, 3 и 4 назад",
         "size": 15},
    ], body_box=(0.75, 4.35, 8.5, 0.6),
        picture="s_window.png", pic_box=(1.5, 1.55, 7.0))

    # ------------------------------------------------------------- модель
    add(L, "Модель", stat("45 с", "обучение на процессоре") + [
        {"t": "LSTM: 2 слоя по 64, 51 651 параметр", "size": 14, "space": 12},
        {"t": "Плато около 0.75 — есть куда расти", "size": 14},
    ], body_box=(6.0, 1.7, 3.3, 2.6),
        picture="s_training.png", pic_box=(0.65, 1.9, 5.1))

    # ------------------------------------------------------ метрика событий
    add(L, "Считаем события, а не точки", [
        {"t": "Угадала 3 интервала из 5 — событие всё равно найдено. "
              "Одно срабатывание на пустом месте — ложная тревога.", "size": 14},
        {"t": "Допуск ±5 минут — это один интервал.", "size": 14, "space": 8},
    ], body_box=(0.8, 4.05, 8.4, 1.0),
        picture="s_event_metric.png", pic_box=(2.0, 1.4, 6.0))

    # ---------------------------------------------------------- результаты
    add(L, "Результаты", stat("0.786", "событийный F1") + [
        {"t": "precision 0.988 · recall 0.653", "size": 14, "space": 12},
        {"t": "81 событие найдено, 43 пропущено", "size": 14},
        {"t": "91 интервал уехал в класс 2", "size": 14},
    ], body_box=(4.4, 1.6, 4.8, 3.0),
        picture="s_confusion.png", pic_box=(0.8, 1.5, 3.2))

    # ------------------------------------------------------ против порога
    add(L, "Модель против порога", [
        {"t": "У порога precision около 0.22: он находит заправки вместе "
              "с колебаниями уровня", "size": 15},
        {"t": "Модель смотрит на форму окна, а не на одно число",
         "size": 15, "space": 10},
    ], body_box=(5.9, 2.0, 3.4, 2.2),
        picture="s_compare.png", pic_box=(0.7, 1.9, 4.9))

    # ---------------------------------------------------------- задание
    add(L, "Задание", [
        {"t": "Поднять событийный F1 на тестовых машинах", "size": 17, "bold": True},
        {"t": "Признаки · нормализация · длина окна · веса классов · "
              "архитектура · сглаживание · порог вероятности · размер интервала",
         "size": 14, "space": 10},
        {"t": "Разбиение не менять, на тестовых машинах не обучаться",
         "size": 14, "space": 10},
    ], body_box=(0.85, 3.5, 8.4, 1.5),
        picture="s_goal.png", pic_box=(1.9, 1.5, 6.2))

    # ------------------------------------------------- вопрос со звёздочкой
    add(L, "Вопрос со звёздочкой", [
        {"t": "Модель для заправок с выключенным двигателем даёт accuracy 100% и F1 = 1.0",
         "size": 15},
        {"t": "Почему это ничего не говорит о качестве детекции?",
         "size": 17, "bold": True, "color": PURPLE, "space": 12},
    ], body_box=(0.85, 4.05, 8.4, 1.0),
        picture="s_leak.png", pic_box=(2.1, 1.5, 5.8))
