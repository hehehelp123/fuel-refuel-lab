# -*- coding: utf-8 -*-
"""Проверка презентации: текст по слайдам, выход за границы, пропорции картинок."""
import glob
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

DECK = "C:/Users/79101/Downloads/Лаба_заправки/Лаба_заправки.pptx"
IMG = "C:/Users/79101/Downloads/Лаба_заправки/slides_img"

print("=== пропорции исходных PNG ===")
for p in sorted(glob.glob(os.path.join(IMG, "*.png"))):
    w, h = Image.open(p).size
    print(f"  {os.path.basename(p):18s} {w}x{h}  aspect {w / h:.3f}")

prs = Presentation(DECK)
SW, SH = prs.slide_width / 914400, prs.slide_height / 914400
print(f"\nхолст: {SW:.2f} x {SH:.2f} дюйма, слайдов: {len(prs.slides)}\n")

problems = []
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for sh in slide.shapes:
        if sh.left is None:
            continue
        x, y = sh.left / 914400, sh.top / 914400
        w = (sh.width or 0) / 914400
        h = (sh.height or 0) / 914400

        if x < 0.25 or y < 0.2 or x + w > SW - 0.25 + 1e-6 or y + h > SH - 0.2 + 1e-6:
            problems.append(f"слайд {i}: {sh.shape_type} выходит за поля "
                            f"({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")

        if sh.shape_type == 13:                      # PICTURE
            iw, ih = sh.image.size
            src, dst = iw / ih, w / h
            if abs(src - dst) / src > 0.02:
                problems.append(f"слайд {i}: картинка искажена, "
                                f"исходная {src:.3f} против {dst:.3f}")

        if sh.has_text_frame and sh.text_frame.text.strip():
            texts.append((sh.text_frame.text.strip(), w, h, sh))

    print(f"--- слайд {i} " + "-" * 46)
    for t, w, h, sh in texts:
        first = t.splitlines()[0]
        print(f"  [{w:4.1f}x{h:4.1f}] {first[:78]}")

    # оценка переполнения: высота считается по каждому абзацу со своим кеглем
    for t, w, h, sh in texts:
        needed = 0.0
        for para in sh.text_frame.paragraphs:
            text = "".join(r.text for r in para.runs)
            size = next((r.font.size.pt for r in para.runs if r.font.size), 16.0)
            char_w = size * 0.5 / 72                  # средняя ширина знака, дюймы
            per_line = max(1, int(w / char_w))
            lines = sum(max(1, -(-len(ln) // per_line))
                        for ln in (text or " ").split("\n"))
            needed += lines * size * 1.25 / 72
            if para.space_before:
                needed += para.space_before.pt / 72
        if needed > h * 1.05:
            problems.append(f'слайд {i}: возможное переполнение '
                            f'({needed:.2f}" против {h:.2f}") — «{t.splitlines()[0][:50]}»')

    # наложение картинок на текстовые блоки
    boxes = []
    for sh in slide.shapes:
        if sh.left is None:
            continue
        rect = (sh.left / 914400, sh.top / 914400,
                (sh.left + (sh.width or 0)) / 914400,
                (sh.top + (sh.height or 0)) / 914400)
        kind = ("pic" if sh.shape_type == 13 else
                ("txt" if sh.has_text_frame and sh.text_frame.text.strip() else None))
        if kind:
            boxes.append((kind, rect, sh))
    for j, (k1, r1, sh1) in enumerate(boxes):
        for k2, r2, sh2 in boxes[j + 1:]:
            if {k1, k2} != {"pic", "txt"}:
                continue
            ox = min(r1[2], r2[2]) - max(r1[0], r2[0])
            oy = min(r1[3], r2[3]) - max(r1[1], r2[1])
            if ox > 0.05 and oy > 0.05:
                txt = (sh1 if k1 == "txt" else sh2).text_frame.text.splitlines()[0][:40]
                problems.append(f'слайд {i}: картинка перекрывает текст '
                                f'({ox:.2f}x{oy:.2f}") — «{txt}»')

print("\n=== замечания ===")
if problems:
    for p in problems:
        print(" ", p)
else:
    print("  чисто")
