# -*- coding: utf-8 -*-
"""Презентация лабораторной на шаблоне ИТМО.

Мастер, макеты и фоны берутся из дипломной презентации. Каждый макет шаблона —
цветная рамка с белым «окном»: фиолетовая, голубая, красная. Обложка — тёмный
фон с логотипом, он задаётся на самом слайде, а не в макете.
"""
import copy
import os
import tempfile
import zipfile

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt

TEMPLATE = "C:/Users/79101/Desktop/Диплом/Презентация_Халеев_М_Д.pptx"
IMG = "C:/Users/79101/Downloads/Лаба_заправки/slides_img"
OUT = "C:/Users/79101/Downloads/Лаба_заправки/Лаба_заправки.pptx"

PURPLE = RGBColor(0x93, 0x07, 0xFE)
CYAN = RGBColor(0x00, 0x8C, 0xB4)      # затемнён, чтобы читался на белом
RED = RGBColor(0xD9, 0x2B, 0x22)
INK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

L_TITLE, L_PURPLE, L_CYAN, L_TWO, L_FINAL = 0, 9, 5, 6, 13
TITLE_PT = 24

# белое «окно» макета: x 0.47..9.53, y 1.06..5.08 дюйма
CARD = (0.6, 1.2, 9.4, 5.0)

prs = Presentation(TEMPLATE)


# ------------------------------------------------- достаём обложку из шаблона
def cover_image_path():
    """Картинка, которой залит фон титульного слайда шаблона."""
    slide = list(prs.slides)[0]
    bg = slide.element.find(qn("p:cSld")).find(qn("p:bg"))
    if bg is None:
        return None
    blip = bg.find(".//" + qn("a:blip"))
    if blip is None:
        return None
    rid = blip.get(qn("r:embed"))
    part = slide.part.rels[rid].target_part
    path = os.path.join(tempfile.gettempdir(), "itmo_cover" +
                        os.path.splitext(str(part.partname))[1])
    with open(path, "wb") as f:
        f.write(part.blob)
    return path


COVER = cover_image_path()

# --- запоминаем номер слайда из шаблона
page_no_proto = None
for slide in prs.slides:
    for sh in slide.shapes:
        if not sh.is_placeholder and sh.has_text_frame and sh.text_frame.text.strip().isdigit():
            page_no_proto = copy.deepcopy(sh._element)
            break
    if page_no_proto is not None:
        break

# --- чистим слайды шаблона, мастер и тема остаются
sld_id_lst = prs.slides._sldIdLst
for sld_id in list(sld_id_lst):
    prs.part.drop_rel(sld_id.rId)
    sld_id_lst.remove(sld_id)


def set_background(slide, image_path):
    image_part, rid = slide.part.get_or_add_image_part(image_path)
    bg = parse_xml(
        f'<p:bg {nsdecls("p", "a", "r")}><p:bgPr>'
        f'<a:blipFill dpi="0" rotWithShape="1"><a:blip r:embed="{rid}"><a:lum/></a:blip>'
        f'<a:srcRect/><a:stretch><a:fillRect/></a:stretch></a:blipFill>'
        f'<a:effectLst/></p:bgPr></p:bg>')
    slide.element.find(qn("p:cSld")).insert(0, bg)


def fill(tf, items):
    """items: строка либо dict(t=текст, size=, color=, bold=, level=, space=)."""
    tf.clear()
    for i, item in enumerate(items):
        d = {"t": item} if isinstance(item, str) else dict(item)
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = d.get("level", 0)
        if d.get("space"):
            para.space_before = Pt(d["space"])
        run = para.add_run()
        run.text = d["t"]
        if d.get("size"):
            run.font.size = Pt(d["size"])
        if d.get("color") is not None:
            run.font.color.rgb = d["color"]
        if d.get("bold"):
            run.font.bold = True


def ph_by_idx(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


FRAMED = None   # заполняется ниже: макеты с цветной рамкой


def new_slide(layout_idx, title, title_color=None):
    slide = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[layout_idx])
    if title_color is None and layout_idx in (L_PURPLE, L_CYAN, L_TWO, L_FINAL):
        title_color = WHITE          # заголовок стоит на цветной рамке
    if slide.shapes.title is not None and title is not None:
        tf = slide.shapes.title.text_frame
        tf.clear()
        run = tf.paragraphs[0].add_run()
        run.text = title
        run.font.size = Pt(TITLE_PT)
        if title_color is not None:
            run.font.color.rgb = title_color
    return slide


def stamp_page_no(slide):
    if page_no_proto is None:
        return
    el = copy.deepcopy(page_no_proto)
    slide.shapes._spTree.append(el)
    for sh in slide.shapes:
        if sh._element is el and sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    run.text = str(len(prs.slides._sldIdLst))
                break


def add(layout_idx, title, items, body_box=None, picture=None, pic_box=None):
    slide = new_slide(layout_idx, title)
    body = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx != 0:
            body = ph
            break
    if body is not None:
        if items:
            fill(body.text_frame, items)
            if body_box:
                body.left, body.top = Inches(body_box[0]), Inches(body_box[1])
                body.width, body.height = Inches(body_box[2]), Inches(body_box[3])
        else:
            body._element.getparent().remove(body._element)
    if picture:
        path = os.path.join(IMG, picture)
        iw, ih = Image.open(path).size
        x, y, w = pic_box
        slide.shapes.add_picture(path, Inches(x), Inches(y),
                                 width=Inches(w), height=Inches(w * ih / iw))
    stamp_page_no(slide)
    return slide


def add_two(title, left, right, layout_idx=None):
    """Две колонки внутри белого поля макета."""
    slide = new_slide(layout_idx or L_PURPLE, title)
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 0:
            ph._element.getparent().remove(ph._element)

    accent = CYAN if (layout_idx or L_PURPLE) == L_CYAN else PURPLE
    for (header, lines), x in ((left, 0.75), (right, 5.15)):
        box = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(4.0), Inches(3.0))
        box.text_frame.word_wrap = True
        fill(box.text_frame,
             [{"t": header, "size": 19, "bold": True, "color": accent}] +
             [{"t": t, "size": 15, "space": 8} for t in lines])
    stamp_page_no(slide)
    return slide


def big(text, color=PURPLE, size=30):
    return {"t": text, "size": size, "color": color, "bold": True, "space": 14}


# ================================================================== титул
slide = new_slide(L_TITLE, None)
for ph in list(slide.placeholders):
    ph._element.getparent().remove(ph._element)
if COVER:
    set_background(slide, COVER)

box = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(8.2), Inches(1.2))
fill(box.text_frame, [{"t": "Детекция заправок\nпо датчику уровня топлива",
                       "size": 30, "color": WHITE, "bold": True}])
box.text_frame.word_wrap = True

box = slide.shapes.add_textbox(Inches(0.9), Inches(4.05), Inches(8.2), Inches(0.9))
fill(box.text_frame, [
    {"t": "Лабораторная работа", "size": 16, "color": WHITE},
    {"t": "50 машин · 6.1 млн сообщений телематики · июль — декабрь 2025",
     "size": 13, "color": RGBColor(0xC9, 0xC9, 0xC9)},
])

# ================================================================== задача
add(L_PURPLE, "Задача", [
    "На машине стоит датчик уровня топлива. Телематический блок шлёт "
    "показания раз в несколько секунд.",
    {"t": "Нужно по этому ряду находить моменты заправок.", "space": 10},
    big("1249 заправок за полгода на 50 машинах"),
])

add(L_CYAN, "Что приходит с машины", [
    "Уровень топлива, зажигание, скорость, время",
    "Сообщения неравномерные: от секунды на ходу до нескольких минут на стоянке",
    "Уровень — в единицах датчика, не в литрах",
    {"t": "Тарировка: у каждой машины своя таблица «показание — литры». "
          "Без неё одна и та же заправка на разных машинах даёт разную амплитуду.",
     "space": 10},
])

add(L_PURPLE, "Порог на приросте не работает", [
    "Датчик запитан от зажигания: на стоянке показаний нет",
    "На ходу уровень гуляет на несколько литров от разгонов и уклонов",
    "Низкий порог ловит эти колебания, высокий пропускает мелкие заправки",
    big("Лучший порог даёт F1 = 0.37"),
])

add_two("Зажигание делит задачу надвое",
        ("Двигатель заглушен", [
            "Роста уровня нет совсем",
            "Виден уровень до остановки и после запуска",
            "711 событий",
        ]),
        ("Двигатель работает", [
            "Уровень растёт монотонно 10–20 минут",
            "Рост виден напрямую",
            "538 событий — с ними и работаем",
        ]))

add(L_CYAN, "Так выглядит заправка", [
    {"t": "Один интервал — 5 минут", "size": 14},
    {"t": "Заправка занимает 3–4 интервала: с 17 до 75 литров", "size": 14},
    {"t": "Серым — разметка оператора, границы приблизительные", "size": 14},
], body_box=(5.9, 1.45, 3.4, 1.8),
    picture="s_event.png", pic_box=(0.65, 2.0, 5.0))

add(L_PURPLE, "Датасет", [
    "85 228 строк: одна строка — 5 минут одной машины",
    "Столбцы: vehicle, ts, fuel, ign, speed, n_msg, label, split",
    "Классы: 0 — ничего, 1 — заправка на ходу, 2 — на стоянке",
    {"t": "Разбиение по машинам: 45 в обучении, 5 в тесте. По точкам делить "
          "нельзя — окна одного события попадут в обе части.", "space": 10},
])

add(L_CYAN, "Целевой класс редкий", [
    {"t": "Модель «всегда ноль» даёт 98% accuracy и не находит "
          "ни одной заправки", "size": 14},
    {"t": "Поэтому в потерях — веса классов: корень из отношения частот",
     "size": 14},
    {"t": "1.8%", "size": 34, "color": PURPLE, "bold": True, "space": 14},
    {"t": "интервалов — заправка на ходу", "size": 13},
], body_box=(5.9, 1.45, 3.4, 2.6),
    picture="s_balance.png", pic_box=(0.65, 1.9, 4.9))

add(L_PURPLE, "Признаки и окна", [
    "Признаки интервала: уровень и четыре разности — на 1, 2, 3 и 4 назад",
    "Короткая разность ловит скачок, длинная — медленный рост",
    "Окно: 7 интервалов подряд, это 35 минут",
    "Класс окна — класс центрального интервала",
    big("74 993 окна в обучении, 9 935 в тесте", size=22),
])

add(L_CYAN, "Модель", [
    "Двухслойная LSTM на 64 нейрона, dropout 0.2, линейный слой на 3 класса",
    "Adam, шаг 0.001, батч 512, 30 эпох",
    "Та же архитектура работает в продуктивном сервисе",
    big("51 651 параметр · 45 секунд обучения на CPU", size=22),
])

add_two("Точечная метрика обманывает", layout_idx=L_CYAN, left=
        ("Угадала 3 интервала из 5", [
            "Точечная метрика: средне",
            "Событийная: заправка найдена",
        ]),
        right=("Одно срабатывание на стоянке", [
            "Точечная метрика: почти незаметно",
            "Событийная: ложная тревога",
        ]))

add(L_PURPLE, "Как считаем события", [
    "Событие — непрерывный блок класса 1 внутри одной машины",
    "Предсказание засчитывается при пересечении с эталоном ±5 минут",
    "Допуск в 5 минут — это один интервал, соседний не считается ошибкой",
    "Попадание на класс 2 не штрафуется: разметке там нельзя доверять",
])

add(L_CYAN, "Результаты", [
    {"t": "0.853", "size": 34, "color": PURPLE, "bold": True},
    {"t": "событийный F1", "size": 13},
    {"t": "precision 1.000 · recall 0.744", "size": 14, "space": 10},
    {"t": "Найдено 96 событий, пропущено 33", "size": 14},
    {"t": "72 интервала уехали в класс 2 — здесь запас", "size": 14},
], body_box=(4.3, 1.5, 4.9, 3.2),
    picture="s_confusion.png", pic_box=(0.7, 1.45, 3.2))

add(L_PURPLE, "Модель против порога", [
    {"t": "Порог перебран от 1 до 20 литров, лучший результат 0.373", "size": 14},
    {"t": "При любом пороге precision около 0.22: правило находит заправки "
          "вместе с колебаниями", "size": 14},
    {"t": "Модель смотрит на форму окна, а не на одно число", "size": 14},
], body_box=(5.7, 1.6, 3.6, 2.4),
    picture="s_compare.png", pic_box=(0.65, 1.9, 4.7))

add(L_CYAN, "Задание", [
    {"t": "Поднять событийный F1 на тестовых машинах", "size": 18, "bold": True},
    {"t": "Планка «хорошо» — выше 0.88", "size": 15, "space": 6},
    {"t": "Признаки · нормализация · длина окна · веса классов · "
          "архитектура · сглаживание предсказаний · порог вероятности · "
          "размер интервала", "size": 14, "space": 12},
    {"t": "Разбиение не менять, на тестовых машинах не обучаться, "
          "гиперпараметры подбирать на обучающей выборке", "size": 14, "space": 10},
])

add(L_FINAL, "Вопрос со звёздочкой", [
    "Вторая модель — для заправок на стоянке — даёт accuracy 100% и F1 = 1.0",
    {"t": "Метка ставится по правилу: три точки, в середине ноль, по краям "
          "уровни стабильны и различаются. Признаки модели — уровень "
          "и его разности.", "space": 10},
    {"t": "Почему это число ничего не говорит о качестве детекции?",
     "size": 18, "bold": True, "color": RED, "space": 14},
])

prs.save(OUT)
print(f"сохранено: {OUT}")
print(f"слайдов: {len(prs.slides._sldIdLst)}")
print(f"обложка: {COVER}")
